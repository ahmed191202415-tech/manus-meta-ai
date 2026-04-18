from typing import Dict, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from app.config import EXPORT_DIR
from app.schemas.report_requests import SavePptxReportRequest
from app.utils.files import sanitize_file_name


def save_pptx_report_local(body: SavePptxReportRequest) -> Dict[str, Any]:
    file_name = sanitize_file_name(body.file_name, "presentation", "pptx")
    output_path = EXPORT_DIR / file_name

    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = body.title
    if body.subtitle:
        slide.placeholders[1].text = body.subtitle

    for spec in body.slides:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = spec.title

        if spec.bullets:
            tx_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(5.2), Inches(4.5))
            tf = tx_box.text_frame
            for i, bullet in enumerate(spec.bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = bullet
                p.font.size = Pt(20)
                p.alignment = PP_ALIGN.RIGHT

        if spec.table_headers:
            rows = len(spec.table_rows) + 1
            cols = len(spec.table_headers)
            table_shape = slide.shapes.add_table(rows, cols, Inches(6.0), Inches(1.5), Inches(3.0), Inches(3.6))
            table = table_shape.table

            for c, header in enumerate(spec.table_headers):
                table.cell(0, c).text = str(header)

            for r, row in enumerate(spec.table_rows, start=1):
                for c, value in enumerate(row):
                    table.cell(r, c).text = str(value)

    prs.save(str(output_path))

    return {
        "success": True,
        "file_name": file_name,
        "file_path": str(output_path),
        "export_dir": str(EXPORT_DIR),
    }
